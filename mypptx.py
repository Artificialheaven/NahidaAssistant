from pptx import Presentation
import json

from datetime import date

today = date.today()


def analyze_ppt_structure(ppt_path):
    prs = Presentation(ppt_path)

    # 整体结构数据
    structure = {
        "slides": [],
        "slide_count": len(prs.slides)
    }

    # 遍历每一页幻灯片
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "layout_type": slide.slide_layout.name,
            "title": None,
            "shapes": []
        }

        # 提取标题
        title_shape = next((s for s in slide.shapes if s.has_text_frame and s.text.strip()), None)
        if title_shape:
            slide_info["title"] = title_shape.text.strip()

        # 遍历所有形状
        for shape in slide.shapes:
            shape_info = {
                "id": shape.shape_id,
                "type": str(shape.shape_type).split(".")[-1],
                "text": [],
                "is_placeholder": shape.is_placeholder
            }

            # 提取文字内容
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        shape_info["text"].append(text)

            # 占位符详细信息
            if shape.is_placeholder:
                shape_info["placeholder_type"] = str(shape.placeholder_format.type).split(".")[-1]

            slide_info["shapes"].append(shape_info)

        structure["slides"].append(slide_info)

    return structure


def get_pptx_tree(filename):
    ppt_struct = analyze_ppt_structure(filename)
    print(f"PPT总页数：{ppt_struct['slide_count']}")

    js = {}
    js["页数"] = str(ppt_struct['slide_count'])
    js["分页"] = []
    for slide in ppt_struct["slides"]:
        print(f"\nSlide {slide['slide_number']} - 版式：{slide['layout_type']}")
        print(f"标题：{slide['title']}")
        print("包含元素：")

        _js = {"标题": slide['title'], "页码": str(slide['slide_number']), "内容": []}
        for shape in slide["shapes"]:
            if not shape['text']:
                continue
            text = "".join(shape['text']) if shape['text'] else ""
            if text == "202X.X":
                text = f"{today.year}.{today.month}"
            if text == "20XX":
                text = f"{today.year}"
            if text == "AiPPT":
                text = f"纳西妲"
            _js["内容"].append(text)

            type_info = f"{shape['type']}(占位符:{shape['is_placeholder']})"
            text_info = " | ".join(shape['text']) if shape['text'] else "无文字"
            print(f"  - ID:{shape['id']} {type_info}: {text_info}")
        js["分页"].append(_js)

    print(js)

    return json.dumps(js)


# 使用示例
if __name__ == "__main__":
    text = get_pptx_tree("需求挖掘：痛点痒点爽点与营销策略.pptx")
    print("")
